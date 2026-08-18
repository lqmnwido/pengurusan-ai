"""Temporal worker for configured Pengurusan AI agents.

Run this process separately from the FastAPI web server:
    python -m open_webui.agents.worker

Custom Python execution is intentionally disabled unless
AGENT_PYTHON_EXECUTION_ENABLED=true. In production this worker should run in a
container with a read-only root filesystem, no host mounts, and network access
disabled unless an agent explicitly requires it.
"""

import asyncio
import io
import json
import os
import subprocess
import sys
import tempfile
from pathlib import Path
from uuid import uuid4

from open_webui.agents import openclaw
from open_webui.agents.agent_graph import build_agent_chain_plan
from open_webui.agents.runtime import openclaw_text
from open_webui.agents.voice_graph import build_voice_flow_plan
from open_webui.models.agents import AgentConfigurations
from open_webui.routers.audio import set_faster_whisper_model
from open_webui.storage.provider import Storage
from pengurusan_temporal_workflows import (
    AgenticChainWorkflow,
    OpenClawChatWorkflow,
    PengurusanAgentWorkflow,
    VoiceIntelligenceWorkflow,
)
from temporalio import activity
from temporalio.client import Client
from temporalio.worker import Worker

PYTHON_HARNESS = r'''
import asyncio
import importlib.util
import inspect
import json
import sys

class Context:
    def log(self, message):
        print(str(message), file=sys.stderr)

spec = importlib.util.spec_from_file_location("uploaded_agent", sys.argv[1])
module = importlib.util.module_from_spec(spec)
spec.loader.exec_module(module)
entrypoint = getattr(module, sys.argv[2])
payload = json.loads(sys.stdin.read())
result = entrypoint(Context(), payload)
if inspect.isawaitable(result):
    result = asyncio.run(result)
print(json.dumps(result))
'''

_WHISPER_LOCK = asyncio.Lock()
_PYANNOTE_LOCK = asyncio.Lock()
_PYANNOTE_PIPELINES: dict[str, object] = {}


async def _await_with_heartbeat(awaitable, details: dict):
    task = asyncio.create_task(awaitable)
    while not task.done():
        await asyncio.wait({task}, timeout=20)
        if not task.done():
            activity.heartbeat(details)
    return await task


def _resource_limits(memory_mb: int):
    def apply():
        try:
            import resource

            memory = memory_mb * 1024 * 1024
            resource.setrlimit(resource.RLIMIT_AS, (memory, memory))
            resource.setrlimit(resource.RLIMIT_CPU, (60, 60))
            resource.setrlimit(resource.RLIMIT_NOFILE, (64, 64))
        except (ImportError, OSError, ValueError):
            pass

    return apply


def run_uploaded_python(code_path: str, entrypoint: str, payload: dict, memory_mb: int, timeout: int):
    local_path = Storage.get_file(code_path)
    with tempfile.TemporaryDirectory(prefix='pengurusan-agent-') as workdir:
        result = subprocess.run(
            [sys.executable, '-I', '-c', PYTHON_HARNESS, local_path, entrypoint],
            input=json.dumps(payload),
            capture_output=True,
            text=True,
            timeout=min(timeout, 3600),
            cwd=workdir,
            env={'PATH': os.getenv('PATH', ''), 'PYTHONIOENCODING': 'utf-8'},
            preexec_fn=_resource_limits(memory_mb) if os.name == 'posix' else None,
            check=False,
        )
    if result.returncode != 0:
        raise RuntimeError(result.stderr.strip() or f'Agent process exited with code {result.returncode}')
    try:
        return json.loads(result.stdout)
    except json.JSONDecodeError as exc:
        raise RuntimeError('Agent output must be one valid JSON value') from exc


def _store_output(agent_id: str, suffix: str, contents: bytes, content_type: str):
    object_name = f'agent-{agent_id}-{uuid4().hex}-{suffix}'
    _, path = Storage.upload_file(
        io.BytesIO(contents),
        object_name,
        {'resource': 'agent-output', 'agent_id': agent_id, 'content_type': content_type},
    )
    return path


def _load_artifact(path: str | None) -> dict:
    if not path:
        return {}
    with open(Storage.get_file(path), encoding='utf-8') as artifact_file:
        return json.load(artifact_file)


def _store_voice_artifact(agent_id: str, job_id: str, stage: str, artifact: dict) -> str:
    safe_job_id = ''.join(character for character in job_id if character.isalnum() or character in '-_')[:80]
    return _store_output(
        agent_id,
        f'{safe_job_id}-{stage}.json',
        json.dumps(artifact, ensure_ascii=False).encode('utf-8'),
        'application/json',
    )


def _assign_speakers(segments: list[dict], speakers: list[dict]) -> list[dict]:
    assigned = []
    for segment in segments:
        best = max(
            speakers,
            key=lambda speaker: max(
                0,
                min(float(segment['end']), float(speaker['end']))
                - max(float(segment['start']), float(speaker['start'])),
            ),
            default=None,
        )
        overlap = (
            max(
                0,
                min(float(segment['end']), float(best['end'])) - max(float(segment['start']), float(best['start'])),
            )
            if best
            else 0
        )
        assigned.append({**segment, 'speaker': best['speaker'] if best and overlap > 0 else 'UNKNOWN'})
    return assigned


async def _analyse_with_openclaw(agent, stage: dict, artifact: dict) -> str:
    transcript = artifact.get('transcript', '')
    chunks = artifact.get('chunks') or _chunk_text(transcript, 12000, 500)
    instruction = stage.get('prompt') or 'Analisis transkrip ini dengan tepat.'
    model_id = stage.get('model_id') or agent.config.model.model_id
    session_prefix = f'voice-{artifact["job_id"]}-{stage["component"]}'

    async def invoke(text: str, session_key: str) -> str:
        response = await _await_with_heartbeat(
            openclaw.invoke_agent(
                agent.config.openclaw.agent_id,
                f'{instruction}\n\nTRANSKRIP:\n{text}',
                model_id,
                timeout=agent.config.orchestration.timeout_seconds,
                session_key=session_key,
            ),
            {'component': stage['component'], 'status': 'analysing'},
        )
        return openclaw_text(response)

    if len(transcript) <= 24000:
        return await invoke(transcript, session_prefix)

    partials = []
    for index, chunk in enumerate(chunks):
        activity.heartbeat({'component': stage['component'], 'chunk': index + 1, 'total': len(chunks)})
        partials.append(await invoke(chunk, f'{session_prefix}-part-{index + 1}'))
    response = await _await_with_heartbeat(
        openclaw.invoke_agent(
            agent.config.openclaw.agent_id,
            f'{instruction}\nGabungkan analisis separa berikut tanpa mengulangi isi atau mereka fakta.\n\n'
            + '\n\n'.join(partials),
            model_id,
            timeout=agent.config.orchestration.timeout_seconds,
            session_key=f'{session_prefix}-final',
        ),
        {'component': stage['component'], 'status': 'synthesising'},
    )
    return openclaw_text(response)


def _chunk_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        if end < len(text):
            boundary = max(text.rfind('\n', start, end), text.rfind('. ', start, end))
            if boundary > start + chunk_size // 2:
                end = boundary + 1
        chunks.append(text[start:end].strip())
        if end >= len(text):
            break
        start = max(start + 1, end - overlap)
    return [chunk for chunk in chunks if chunk]


def _extract_document(path: str) -> str:
    extension = Path(path).suffix.lower()
    if extension in {'.txt', '.md', '.csv', '.json'}:
        return Path(path).read_text(encoding='utf-8', errors='replace')
    if extension == '.pdf':
        from pypdf import PdfReader

        return '\n\n'.join(page.extract_text() or '' for page in PdfReader(path).pages)
    if extension == '.docx':
        from docx import Document

        return '\n'.join(paragraph.text for paragraph in Document(path).paragraphs)
    if extension == '.pptx':
        from pptx import Presentation

        return '\n'.join(
            shape.text for slide in Presentation(path).slides for shape in slide.shapes if hasattr(shape, 'text')
        )
    raise RuntimeError(f'Unsupported document type: {extension or "unknown"}')


def run_builtin_agent(agent, payload: dict):
    file_path = payload.get('file_path')
    if not file_path:
        raise RuntimeError('Built-in file agents require payload.file_path from the platform file store')
    local_path = Storage.get_file(file_path)

    if 'image_recognition' in agent.config.capabilities or 'ocr' in agent.config.capabilities:
        from PIL import Image

        image = Image.open(local_path).convert('RGB')
        result = {'status': 'completed'}
        if 'image_recognition' in agent.config.capabilities:
            from transformers import AutoModelForMultimodalLM, AutoProcessor

            model_id = agent.config.model.model_id or 'Salesforce/blip-image-captioning-base'
            processor = AutoProcessor.from_pretrained(model_id)
            model = AutoModelForMultimodalLM.from_pretrained(model_id)
            inputs = processor(images=image, return_tensors='pt')
            output = model.generate(**inputs, max_new_tokens=min(agent.config.model.max_tokens, 256))
            result['caption'] = processor.batch_decode(output, skip_special_tokens=True)[0].strip()
        if 'ocr' in agent.config.capabilities:
            from rapidocr_onnxruntime import RapidOCR

            ocr_result, _ = RapidOCR()(local_path)
            result['ocr_text'] = '\n'.join(item[1] for item in (ocr_result or []))
        if agent.config.storage.save_outputs:
            result['output_path'] = _store_output(
                agent.id, 'image-analysis.json', json.dumps(result).encode(), 'application/json'
            )
        return result

    if agent.config.template == 'voice-transcription' or 'voice_to_text' in agent.config.capabilities:
        model = set_faster_whisper_model(agent.config.model.stt_model, True)
        segments, info = model.transcribe(local_path, vad_filter=True)
        transcript = ' '.join(segment.text.strip() for segment in segments).strip()
        result = {
            'status': 'completed',
            'language': getattr(info, 'language', None),
            'duration': getattr(info, 'duration', None),
            'transcript': transcript,
        }
        if 'diarization' in agent.config.capabilities:
            if 'speechbrain-diarization' in agent.config.components:
                raise RuntimeError(
                    'SpeechBrain + Silero diarization uses the AsrWorkflow::start orchestration profile and ASR worker queues.'
                )
            try:
                from pyannote.audio import Pipeline

                token = os.getenv('HF_TOKEN') or os.getenv('HUGGINGFACE_TOKEN')
                pipeline = Pipeline.from_pretrained('pyannote/speaker-diarization-community-1', token=token)
                diarized = pipeline(local_path)
                annotation = getattr(diarized, 'exclusive_speaker_diarization', None) or getattr(
                    diarized, 'speaker_diarization', diarized
                )
                speakers = []
                if hasattr(annotation, 'itertracks'):
                    speakers = [
                        {'start': turn.start, 'end': turn.end, 'speaker': speaker}
                        for turn, _, speaker in annotation.itertracks(yield_label=True)
                    ]
                else:
                    speakers = [
                        {'start': turn.start, 'end': turn.end, 'speaker': speaker} for turn, speaker in annotation
                    ]
                result['speakers'] = speakers
            except Exception as exc:
                raise RuntimeError(
                    f'Diarization failed. Verify pyannote installation, model terms, and HF_TOKEN: {exc}'
                )
        chunks = []
        if agent.config.chunking.enabled:
            chunks = _chunk_text(
                transcript,
                agent.config.chunking.chunk_size,
                agent.config.chunking.chunk_overlap,
            )
            result['chunk_count'] = len(chunks)
        if agent.config.storage.save_outputs:
            artifact = {**result, 'chunks': chunks}
            result['output_path'] = _store_output(
                agent.id, 'transcript.json', json.dumps(artifact).encode(), 'application/json'
            )
        elif chunks:
            result['chunks'] = chunks
        return result

    text = _extract_document(local_path)
    chunks = (
        _chunk_text(text, agent.config.chunking.chunk_size, agent.config.chunking.chunk_overlap)
        if agent.config.chunking.enabled
        else [text]
    )
    result = {
        'status': 'completed',
        'characters': len(text),
        'chunk_count': len(chunks),
        'chunking_strategy': agent.config.chunking.strategy,
    }
    if agent.config.storage.save_outputs:
        artifact = {**result, 'chunks': chunks}
        result['output_path'] = _store_output(
            agent.id, 'chunks.json', json.dumps(artifact).encode(), 'application/json'
        )
        result['preview'] = chunks[:3]
    else:
        result['chunks'] = chunks
    return result


@activity.defn(name='pengurusan_ai.execute_agent')
async def execute_agent(input_data: dict):
    agent_id = input_data['agent_id']
    payload = input_data.get('payload', {})
    agent = await AgentConfigurations.get(agent_id)
    if not agent or not agent.is_active:
        raise RuntimeError('Agent is missing or inactive')
    if agent.config.runtime.mode == 'python':
        if os.getenv('AGENT_PYTHON_EXECUTION_ENABLED', 'false').lower() != 'true':
            raise RuntimeError('Set AGENT_PYTHON_EXECUTION_ENABLED=true only inside an isolated worker container.')
        if not agent.code_path:
            raise RuntimeError('No Python package has been uploaded for this agent')
        if agent.config.runtime.allow_network and os.getenv('AGENT_WORKER_NETWORK_ENABLED', 'false').lower() != 'true':
            raise RuntimeError('This agent requests network access but the worker has not enabled it')

        return await asyncio.to_thread(
            run_uploaded_python,
            agent.code_path,
            agent.config.runtime.entrypoint,
            payload,
            agent.config.runtime.memory_mb,
            agent.config.orchestration.timeout_seconds,
        )

    return await asyncio.to_thread(run_builtin_agent, agent, payload)


@activity.defn(name='pengurusan_ai.openclaw.invoke')
async def invoke_openclaw_agent(input_data: dict):
    """Invoke OpenClaw from a Temporal activity, never from workflow code."""
    agent = await AgentConfigurations.get(input_data['agent_id'])
    if not agent or not agent.is_active:
        raise RuntimeError('Agent is missing or inactive')
    if agent.config.framework != 'openclaw' or not agent.config.openclaw.agent_id:
        raise RuntimeError('Agent is not linked to OpenClaw')
    if not openclaw.enabled():
        raise RuntimeError('OpenClaw is not enabled')

    return await openclaw.invoke_agent(
        agent.config.openclaw.agent_id,
        input_data['message'],
        agent.config.model.model_id,
        timeout=int(input_data.get('timeout_seconds', 900)),
        session_key=input_data.get('session_key'),
    )


@activity.defn(name='pengurusan_ai.langgraph.voice.plan')
async def plan_voice_flow(input_data: dict):
    return await asyncio.to_thread(build_voice_flow_plan, input_data['flow'])


@activity.defn(name='pengurusan_ai.langgraph.agent.plan')
async def plan_agent_chain(input_data: dict):
    return await asyncio.to_thread(build_agent_chain_plan, input_data['nodes'])


@activity.defn(name='pengurusan_ai.openclaw.chain.node')
async def execute_agentic_chain_node(input_data: dict):
    node = input_data['node']
    agent = await AgentConfigurations.get(node['agent_id'])
    if not agent or not agent.is_active or not agent.config.openclaw.agent_id:
        raise RuntimeError('Workflow node agent is missing, inactive, or not linked to OpenClaw')
    instruction = node.get('instruction', '').strip()
    message = build_agentic_node_message(input_data, instruction)
    result = await openclaw.invoke_agent(
        agent.config.openclaw.agent_id,
        message,
        node.get('model_id') or agent.config.model.model_id,
        timeout=int(input_data.get('timeout_seconds', 900)),
        session_key=f'{input_data["session_key"]}-node-{input_data["node_index"] + 1}',
    )
    return {
        'agent_id': agent.id,
        'agent_name': agent.name,
        'openclaw_agent_id': agent.config.openclaw.agent_id,
        'model_id': node.get('model_id') or agent.config.model.model_id,
        'instruction': instruction
        or (
            'Memproses permintaan pengguna'
            if input_data['node_index'] == 0
            else 'Memproses output daripada agent sebelumnya'
        ),
        'text': openclaw_text(result),
    }


def build_agentic_node_message(input_data: dict, instruction: str = '') -> str:
    """Give each node only the document it is responsible for processing."""
    if input_data.get('node_index', 0) == 0:
        message = input_data['original_message']
    else:
        message = (
            'Process only the complete document inside <source_document> according to your '
            'agent identity and the instruction for this workflow step. Do not use or infer the '
            'original user request. Preserve the source document unless your assigned task '
            'explicitly requires transforming it. Return only the completed result.\n\n'
            f'<source_document>\n{input_data["previous_output"]}\n</source_document>'
        )
    return f'{instruction}\n\n{message}' if instruction else message


@activity.defn(name='pengurusan_ai.langgraph.voice.node')
async def execute_voice_node(input_data: dict):
    state = input_data['state']
    stage = input_data['stage']
    component = stage['component']
    agent = await AgentConfigurations.get(stage.get('agent_id') or state['agent_id'])
    if not agent or not agent.is_active:
        raise RuntimeError(f'Node agent for {component} is missing or inactive')

    activity.heartbeat({'component': component, 'status': 'started'})
    artifact = _load_artifact(state.get('artifact_path'))
    artifact.update(
        {
            'agent_id': agent.id,
            'openclaw_agent_id': agent.config.openclaw.agent_id,
            'agentic_workflow_id': state.get('agentic_workflow_id'),
            'job_id': state['job_id'],
            'input_file': state['input_file'],
        }
    )

    if component == 'faster-whisper':
        local_path = Storage.get_file(state['input_file'])

        def transcribe():
            model = set_faster_whisper_model(stage.get('model_id') or 'small', True)
            segments, info = model.transcribe(local_path, vad_filter=True)
            return (
                [{'start': segment.start, 'end': segment.end, 'text': segment.text.strip()} for segment in segments],
                info,
            )

        async with _WHISPER_LOCK:
            segment_list, info = await _await_with_heartbeat(
                asyncio.to_thread(transcribe),
                {'component': component, 'status': 'transcribing'},
            )
        artifact.update(
            {
                'language': getattr(info, 'language', None),
                'duration': getattr(info, 'duration', None),
                'segments': segment_list,
                'transcript': ' '.join(item['text'] for item in segment_list).strip(),
            }
        )
    elif component == 'pyannote-diarization':
        from pyannote.audio import Pipeline

        token = os.getenv('HF_TOKEN') or os.getenv('HUGGINGFACE_TOKEN')
        if not token:
            raise RuntimeError('Pyannote requires HF_TOKEN and accepted Hugging Face model terms')
        local_path = Storage.get_file(state['input_file'])
        model_id = stage.get('model_id') or 'pyannote/speaker-diarization-community-1'
        async with _PYANNOTE_LOCK:
            pipeline = _PYANNOTE_PIPELINES.get(model_id)
            if pipeline is None:
                pipeline = await _await_with_heartbeat(
                    asyncio.to_thread(Pipeline.from_pretrained, model_id, token=token),
                    {'component': component, 'status': 'loading-model'},
                )
                _PYANNOTE_PIPELINES[model_id] = pipeline
            diarized = await _await_with_heartbeat(
                asyncio.to_thread(pipeline, local_path),
                {'component': component, 'status': 'diarizing'},
            )
        annotation = getattr(diarized, 'exclusive_speaker_diarization', None) or getattr(
            diarized, 'speaker_diarization', diarized
        )
        if hasattr(annotation, 'itertracks'):
            speakers = [
                {'start': turn.start, 'end': turn.end, 'speaker': speaker}
                for turn, _, speaker in annotation.itertracks(yield_label=True)
            ]
        else:
            speakers = [{'start': turn.start, 'end': turn.end, 'speaker': speaker} for turn, speaker in annotation]
        artifact['speakers'] = speakers
        artifact['segments'] = _assign_speakers(artifact.get('segments', []), speakers)
    elif component == 'speechbrain-diarization':
        raise RuntimeError(
            'SpeechBrain diarization is provided by the external temporal-worker ASR profile. '
            'Choose Pyannote for this native flow or configure the AsrWorkflow::start profile.'
        )
    elif component == 'transcript-chunking':
        options = stage.get('options') or {}
        chunk_size = int(options.get('chunk_size', agent.config.chunking.chunk_size))
        chunk_overlap = int(options.get('chunk_overlap', agent.config.chunking.chunk_overlap))
        if chunk_size < 100 or chunk_overlap < 0 or chunk_overlap >= chunk_size:
            raise RuntimeError('Chunk size must be at least 100 and overlap must be smaller than chunk size')
        artifact['chunks'] = _chunk_text(
            artifact.get('transcript', ''),
            chunk_size,
            chunk_overlap,
        )
        artifact['chunking'] = {
            'strategy': stage.get('model_id') or agent.config.chunking.strategy,
            'chunk_size': chunk_size,
            'chunk_overlap': chunk_overlap,
        }
    elif component == 'openclaw-topic':
        artifact['topics'] = await _analyse_with_openclaw(agent, stage, artifact)
    elif component == 'openclaw-summary':
        artifact['summary'] = await _analyse_with_openclaw(agent, stage, artifact)
    elif component == 'openclaw-thematic':
        artifact['thematic'] = await _analyse_with_openclaw(agent, stage, artifact)
    else:
        raise RuntimeError(f'Unsupported Voice Intelligence component: {component}')

    artifact.setdefault('completed_components', []).append(component)
    artifact_path = _store_voice_artifact(agent.id, state['job_id'], component, artifact)
    activity.heartbeat({'component': component, 'status': 'completed', 'artifact_path': artifact_path})
    return {**state, 'artifact_path': artifact_path, 'last_component': component}


async def main():
    address = os.getenv('TEMPORAL_ADDRESS', 'localhost:7233')
    namespace = os.getenv('TEMPORAL_NAMESPACE', 'default')
    task_queue = os.getenv('TEMPORAL_TASK_QUEUE', 'pengurusan-ai-agents')
    tls = os.getenv('TEMPORAL_TLS', 'false').lower() == 'true'
    api_key = os.getenv('TEMPORAL_API_KEY', '') or None
    client = await Client.connect(address, namespace=namespace, tls=tls, api_key=api_key)
    worker = Worker(
        client,
        task_queue=task_queue,
        workflows=[PengurusanAgentWorkflow, OpenClawChatWorkflow, AgenticChainWorkflow, VoiceIntelligenceWorkflow],
        activities=[
            execute_agent,
            invoke_openclaw_agent,
            plan_agent_chain,
            execute_agentic_chain_node,
            plan_voice_flow,
            execute_voice_node,
        ],
    )
    print(f'Agent worker listening on Temporal task queue: {task_queue}')
    await worker.run()


if __name__ == '__main__':
    asyncio.run(main())
